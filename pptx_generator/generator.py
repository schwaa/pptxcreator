from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import PP_PLACEHOLDER_TYPE, MSO_SHAPE_TYPE
from pptx.enum.chart import XL_CHART_TYPE
from pptx.chart.data import CategoryChartData
import os
import json
import logging

logging.basicConfig(level=logging.INFO, format='%(levelname)s: %(module)s - %(message)s')

def generate_presentation(data_filepath, template_filepath, output_filepath):
    """
    Generates a PowerPoint presentation based on structured data and a template.
    
    Args:
        data_filepath (str): Path to the JSON file containing presentation data
        template_filepath (str): Path to the template PPTX file
        output_filepath (str): Path where the generated PPTX will be saved
    """
    try:
        with open(data_filepath, 'r', encoding='utf-8') as f:
            data = json.load(f)
    except Exception as e:
        logging.error(f"Error loading data file: {e}")
        raise

    try:
        prs = Presentation(template_filepath)
        if len(prs.slides._sldIdLst) > 0:
            for _ in range(len(prs.slides._sldIdLst)):
                rId = prs.slides._sldIdLst[0].rId
                prs.part.drop_rel(rId)
                del prs.slides._sldIdLst[0]
    except Exception as e:
        logging.error(f"Error loading template or clearing slides: {e}")
        raise

    for slide_idx, slide_plan in enumerate(data.get("slides", [])):
        layout_name = slide_plan.get("layout")
        content_placeholders = slide_plan.get("placeholders", {})
        logging.info(f"Processing slide {slide_idx} with layout: '{layout_name}'")

        slide_layout = None
        for layout_obj in prs.slide_layouts:
            if layout_obj.name == layout_name:
                slide_layout = layout_obj
                break
        if not slide_layout:
            logging.warning(f"  Layout '{layout_name}' not found in template. Skipping slide.")
            continue

        slide = prs.slides.add_slide(slide_layout)
        actual_layout_placeholders = {ph.name.strip(): ph for ph in slide.placeholders}
        logging.info(f"  Layout '{layout_name}' (actual slide instance) has placeholders: {list(actual_layout_placeholders.keys())}")

        for placeholder_name_from_json, value in content_placeholders.items():
            placeholder_name_from_json_stripped = placeholder_name_from_json.strip()
            logging.info(f"  Attempting to fill placeholder from JSON: '{placeholder_name_from_json_stripped}' with value: '{str(value)[:100]}...'")
            
            target_shape = actual_layout_placeholders.get(placeholder_name_from_json_stripped)

            if target_shape:
                ph_format = target_shape.placeholder_format
                logging.info(f"    Found matching placeholder shape named '{target_shape.name.strip()}' (idx: {ph_format.idx if hasattr(ph_format, 'idx') else 'N/A'}, type: {ph_format.type if hasattr(ph_format, 'type') else 'N/A'}).")
                
                # Heuristic to determine if the value is intended as an image
                # This should ideally be more robust, perhaps with a type field in presentation.json
                is_image_value = isinstance(value, str) and \
                                 (value.lower().endswith(('.png', '.jpg', '.jpeg', '.gif', '.bmp', '.tiff')) or \
                                  value.startswith("images/") or \
                                  value.startswith("projects/") or \
                                  "image of" in value.lower() or \
                                  "photo of" in value.lower() or \
                                  "picture of" in value.lower() or \
                                  "graphic of" in value.lower() or \
                                  "diagram of" in value.lower() or \
                                  "screenshot of" in value.lower() or \
                                  "illustration of" in value.lower()
                                 )

                if is_image_value:
                    if ph_format.type == PP_PLACEHOLDER_TYPE.PICTURE:
                        image_path_str = value # The string value from JSON, potentially a path or description
                        
                        # Attempt to resolve image_path_str to an actual file path
                        resolved_image_path = image_path_str 
                        if not os.path.isabs(image_path_str):
                            # Try path relative to project's images/ folder or project root
                            # Assumes data_filepath is like 'projects/project_name/output/presentation.json'
                            project_dir_guess = os.path.dirname(os.path.dirname(data_filepath)) # e.g. projects/robotics
                            
                            # Path patterns to check:
                            # 1. Direct value if it's like "images/my_image.png" or "projects/proj/images/my_image.png"
                            #    (relative to current project or absolute if "projects/" is at root)
                            # 2. Relative to the project's "images" folder: projects/project_name/images/value
                            # 3. Relative to the project's root: projects/project_name/value
                            # 4. Relative to the main pptxcreator "images" folder (less likely for project-specific images)
                            
                            potential_paths = []
                            # Path relative to project's images folder (e.g. projects/robotics/images/my_image.png)
                            potential_paths.append(os.path.join(project_dir_guess, "images", os.path.basename(image_path_str)))
                             # Path relative to project root (e.g. projects/robotics/my_image.png or projects/robotics/images/my_image.png if value includes "images/")
                            potential_paths.append(os.path.join(project_dir_guess, image_path_str))

                            # Check if any potential path exists
                            found_path = False
                            for p_path in potential_paths:
                                if os.path.exists(p_path):
                                    resolved_image_path = p_path
                                    found_path = True
                                    break
                            if not found_path and os.path.exists(image_path_str): # Check original value as a relative path from execution dir
                                resolved_image_path = image_path_str
                                found_path = True
                            
                            if not found_path: # If still not found, log and skip
                                logging.warning(f"      Image file not found. Searched for '{image_path_str}' and derived paths like '{potential_paths}'. Skipping image insertion for Picture Placeholder '{target_shape.name.strip()}'.")
                                continue # Skip to next placeholder

                        # At this point, resolved_image_path should be the one to use if it exists
                        if os.path.exists(resolved_image_path):
                            try:
                                logging.info(f"      Inserting image '{resolved_image_path}' into Picture Placeholder '{target_shape.name.strip()}'")
                                placeholder_shape = target_shape
                                placeholder_shape.insert_picture(resolved_image_path)
                            except Exception as e:
                                logging.warning(f"      Could not insert image '{resolved_image_path}' into '{placeholder_name_from_json_stripped}': {e}")
                        else:
                            logging.warning(f"      Image file '{resolved_image_path}' (derived from '{value}') does not exist. Skipping image insertion for Picture Placeholder '{target_shape.name.strip()}'.")
                    else: # is_image_value is true, but placeholder is not a PICTURE type
                        logging.error(
                            f"      CRITICAL ERROR: Attempted to treat value '{str(value)[:100]}...' as an image "
                            f"for a placeholder that is NOT a PICTURE placeholder (type: {ph_format.type}). "
                            f"Placeholder name: '{target_shape.name.strip()}'. Skipping this placeholder content."
                        )
                
                # Handle text content (if not an image value OR if it's a non-picture placeholder)
                elif hasattr(target_shape, "text_frame"):
                    # Explicitly prevent writing text to picture placeholders if image insertion failed or wasn't appropriate
                    if ph_format.type == PP_PLACEHOLDER_TYPE.PICTURE:
                        logging.warning(f"      Skipping text insertion for value '{str(value)[:100]}...' into Picture Placeholder '{target_shape.name.strip()}'. This placeholder is primarily for images.")
                    else:
                        try:
                            tf = target_shape.text_frame
                            tf.clear() 

                            if isinstance(value, list):
                                logging.info(f"      Setting list content for '{placeholder_name_from_json_stripped}'.")
                                for i, item_text in enumerate(value):
                                    p = tf.add_paragraph()
                                    p.text = str(item_text)
                                    # p.level = 0 # Optional: set bullet level
                            else: 
                                logging.info(f"      Setting string content for '{placeholder_name_from_json_stripped}'.")
                                p = tf.paragraphs[0] if len(tf.paragraphs) > 0 else tf.add_paragraph()
                                p.text = str(value)
                            # tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE # Optional
                        except Exception as e:
                            logging.warning(f"      Could not set text for '{placeholder_name_from_json_stripped}': {e}")
                else:
                    logging.warning(f"      Placeholder '{target_shape.name.strip()}' is not a picture placeholder and has no text_frame. Cannot insert value '{str(value)[:100]}...'.")
            else:
                logging.warning(f"    Placeholder '{placeholder_name_from_json_stripped}' (from JSON) not found in layout '{layout_name}' actual placeholders: {list(actual_layout_placeholders.keys())}.")

    try:
        prs.save(output_filepath)
        print(f"Success! Presentation saved to {output_filepath}")
        return True
    except Exception as e:
        logging.error(f"Error saving presentation: {e}")
        return False
