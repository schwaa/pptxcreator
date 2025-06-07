import re
from pptx.enum.shapes import PP_PLACEHOLDER_TYPE, MSO_SHAPE_TYPE

def clean_text_for_presentation(text: str) -> str:
    """
    Removes leading markdown-style list formatting and surrounding bold markers.
    """
    if not isinstance(text, str):
        # Attempt to convert to string, or return as is if not possible/sensible
        try:
            return str(text) 
        except:
            return text # Or raise an error, or return a default string

    # Remove leading hyphens, asterisks, and whitespace
    # Handles cases like "- item", "* item", "  - item"
    text = re.sub(r'^\s*[-\*\u2022]\s*', '', text) # \u2022 is the bullet character •

    # Remove surrounding bold markers if they exist (e.g., **text**)
    if text.startswith('**') and text.endswith('**') and len(text) > 4:
        text = text[2:-2]
    
    # Remove surrounding italic markers if they exist (e.g., *text* or _text_)
    # Check length to avoid stripping single '*' or '_' if they are part of the text
    if text.startswith('*') and text.endswith('*') and len(text) > 2:
        text = text[1:-1]
    elif text.startswith('_') and text.endswith('_') and len(text) > 2:
        text = text[1:-1]

    return text.strip()

def find_placeholder_by_name(slide, placeholder_name):
    """Find a placeholder by its name in a slide."""
    if not placeholder_name:
        return None
        
    # Try exact match first
    for shape in slide.placeholders:
        if shape.name == placeholder_name:
            return shape
    
    # Try case-insensitive match
    for shape in slide.placeholders:
        if shape.name.lower() == placeholder_name.lower():
            return shape
            
    # Try partial match
    search_term = placeholder_name.lower()
    for shape in slide.placeholders:
        if search_term in shape.name.lower():
            return shape
    
    return None

def find_text_placeholder_by_idx(slide, idx):
    """Find a text placeholder by its index in a slide."""
    for shape in slide.placeholders:
        if shape.placeholder_format.idx == idx:
            return shape
    return None

def find_picture_placeholder_by_type(slide):
    """Find a picture placeholder in a slide."""
    # Try finding by placeholder type
    for shape in slide.placeholders:
        if shape.is_placeholder and shape.placeholder_format.type == PP_PLACEHOLDER_TYPE.PICTURE:
            return shape
    
    # Try finding by name containing "Picture"
    for shape in slide.placeholders:
        if shape.name and "picture" in shape.name.lower():
            return shape
            
    # Try finding by name containing "Image"
    for shape in slide.placeholders:
        if shape.name and "image" in shape.name.lower():
            return shape
            
    return None

def populate_text_placeholder(placeholder_shape, text_runs):
    """
    Populate a text placeholder with content, handling lists of text runs.
    Each item in text_runs can be a string or a tuple (text, bold, italic, font_size, font_name, color_rgb).
    If it's a list, it's assumed to be bullet points.
    """
    if not placeholder_shape or not hasattr(placeholder_shape, 'text_frame'):
        return

    tf = placeholder_shape.text_frame
    tf.clear()  # Clear existing content

    if isinstance(text_runs, str): # Single string
        p = tf.paragraphs[0] if tf.paragraphs else tf.add_paragraph()
        run = p.add_run()
        run.text = text_runs
    elif isinstance(text_runs, list): # List of strings (bullet points) or rich text tuples
        for i, item_text in enumerate(text_runs):
            if not item_text and i == 0 and len(text_runs) == 1: # Handle single empty string in list
                 p = tf.paragraphs[0] if tf.paragraphs else tf.add_paragraph()
                 p.text = "" # Explicitly set to empty
                 continue

            if not item_text and isinstance(item_text, str): # Skip empty strings in a list unless it's the only item
                if len(text_runs) > 1:
                    continue
            
            p = tf.add_paragraph()
            p.text = str(item_text) # Ensure it's a string
            if len(text_runs) > 1: # Apply bullet points if more than one item
                 p.level = 0 # Default bullet level
    else:
        # Fallback for other types, convert to string
        p = tf.paragraphs[0] if tf.paragraphs else tf.add_paragraph()
        run = p.add_run()
        run.text = str(text_runs)


def populate_image_placeholder(placeholder_shape, image_path):
    """Populate an image placeholder with an image."""
    if not placeholder_shape or not image_path:
        return

    try:
        # Basic check if image_path is somewhat valid, could be improved
        if not isinstance(image_path, str) or not os.path.exists(image_path):
            print(f"Warning: Image path '{image_path}' not found or invalid. Attempting to use dummy image.")
            # Ensure images directory exists
            if not os.path.exists("images"):
                os.makedirs("images", exist_ok=True)
            
            # Path for the dummy image
            dummy_image_path = "images/dummy_image.png"

            # Create dummy image if it doesn't exist
            if not os.path.exists(dummy_image_path):
                try:
                    from PIL import Image, ImageDraw
                    img = Image.new('RGB', (800, 600), color='lightblue')
                    draw = ImageDraw.Draw(img)
                    draw.text((10,10), "Image not found", fill=(0,0,0))
                    img.save(dummy_image_path)
                    print(f"Created dummy image at {dummy_image_path}")
                except ImportError:
                    print("PIL (Pillow) is not installed. Cannot create dummy image. Please install Pillow.")
                    return # Cannot proceed without PIL if original image is missing and dummy can't be made
                except Exception as e_pil:
                    print(f"Error creating dummy image: {e_pil}")
                    return # Cannot proceed if dummy image creation fails
            image_path = dummy_image_path
            
        placeholder_shape.insert_picture(image_path)

    except FileNotFoundError: # Should be caught by os.path.exists now, but as a safeguard
        print(f"Error: File not found at '{image_path}' during insert_picture.")
    except Exception as e:
        print(f"Error populating image placeholder: {e}")
        pass  # Silently ignore other image placeholder errors for now
