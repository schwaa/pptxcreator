import json
import os
import argparse
from pptx import Presentation
# from pptx.enum.shapes import MSO_PLACEHOLDER # Import MSO_PLACEHOLDER # This was causing an error

def analyze_template(template_filepath: str, output_filepath: str = None): # output_filepath is now optional
    """
    Analyzes a PowerPoint template.
    If output_filepath is provided, saves a JSON file listing each layout
    and its placeholder details (name, type, position, dimensions).
    Always returns a dictionary of the layout data.
    This version instantiates each layout to get the most accurate placeholder information.

    Args:
        template_filepath (str): Path to the template PPTX file.
        output_filepath (str, optional): Path to save the output layouts.json file. Defaults to None.

    Returns:
        dict: A dictionary containing the layout data.
              Example: {
                          "layouts": [
                            {
                              "name": "Layout Name",
                              "placeholders": {
                                "Placeholder Name 1": {"type": "TITLE", "left_pct": 5, ...},
                                "Placeholder Name 2": {"type": "BODY", "left_pct": 10, ...}
                              }
                            }
                          ]
                       }
                       or {"layouts": [], "error": "Error message"} if loading fails.
    """
    print(f"Analyzing template: {template_filepath}")
    
    try:
        abs_template_filepath = os.path.abspath(template_filepath)
    except Exception as e:
        print(f"Warning: Could not determine absolute path for '{template_filepath}'. Storing as provided. Error: {e}")
        abs_template_filepath = template_filepath

    layout_data = {
        "source_template_path": abs_template_filepath,
        "layouts": []
    }

    try:
        initial_prs = Presentation(template_filepath)
        num_layouts = len(initial_prs.slide_layouts)
        # Get overall slide dimensions from the presentation (they are consistent)
        slide_width_emu = initial_prs.slide_width
        slide_height_emu = initial_prs.slide_height
    except Exception as e:
        error_message = f"Error: Could not load template '{template_filepath}' for initial analysis. Details: {e}"
        print(error_message)
        return {"layouts": [], "error": error_message}

    for i in range(num_layouts):
        try:
            # Reloading prs_temp for each layout is a robust way to ensure clean state,
            # though slide_width/height could be taken from initial_prs.
            # For consistency in accessing slide_layout, we reload.
            prs_temp = Presentation(template_filepath) 
        except Exception as e:
            error_detail = f"Error: Could not re-load template '{template_filepath}' for layout {i}. Details: {e}"
            print(error_detail)
            layout_data["layouts"].append({
                "name": f"Layout index {i} (Load Error)",
                "placeholders": {}, # Changed from list to dict
                "error": error_detail
            })
            continue

        if i >= len(prs_temp.slide_layouts):
            print(f"Error: Layout index {i} out of bounds.")
            continue
            
        slide_layout = prs_temp.slide_layouts[i]
        
        placeholders_details = {}
        
        # Use slide dimensions from the initially loaded presentation for consistency
        # as prs_temp is reloaded, but these properties should be presentation-wide.
        # If prs_temp.slide_width/height were different, it would be strange.
        # Using initial_prs values ensures we use one consistent set.

        source_shapes = []
        try:
            temp_slide_for_analysis = prs_temp.slides.add_slide(slide_layout)
            source_shapes = temp_slide_for_analysis.placeholders 
        except Exception as e:
            print(f"Warning: Could not add temporary slide for layout '{slide_layout.name}'. Error: {e}. Falling back to layout's own placeholders.")
            source_shapes = slide_layout.placeholders

        for shape in source_shapes:
            if not (hasattr(shape, 'is_placeholder') and shape.is_placeholder and \
                    hasattr(shape, 'name') and shape.name):
                print(f"Info: Skipping shape in layout '{slide_layout.name}' as it's not a valid named placeholder. Shape: {shape}")
                continue

            name = shape.name
            placeholder_type_str = "UNKNOWN"
            
            if hasattr(shape, 'placeholder_format') and shape.placeholder_format and \
               hasattr(shape.placeholder_format, 'type') and shape.placeholder_format.type:
                # Convert enum to string name (e.g., MSO_PLACEHOLDER.TITLE -> "TITLE")
                try:
                    placeholder_type_str = shape.placeholder_format.type.name 
                except AttributeError: # Should not happen if type is a valid MSO_PLACEHOLDER enum
                    placeholder_type_str = str(shape.placeholder_format.type) # fallback to raw value if .name fails
            
            left_emu = getattr(shape, 'left', 0)
            top_emu = getattr(shape, 'top', 0)
            width_emu = getattr(shape, 'width', 0)
            height_emu = getattr(shape, 'height', 0)

            left_pct = round((left_emu / slide_width_emu) * 100, 2) if slide_width_emu != 0 else 0
            top_pct = round((top_emu / slide_height_emu) * 100, 2) if slide_height_emu != 0 else 0
            width_pct = round((width_emu / slide_width_emu) * 100, 2) if slide_width_emu != 0 else 0
            height_pct = round((height_emu / slide_height_emu) * 100, 2) if slide_height_emu != 0 else 0
            
            placeholders_details[name] = {
                "type": placeholder_type_str,
                "left_pct": left_pct,
                "top_pct": top_pct,
                "width_pct": width_pct,
                "height_pct": height_pct
            }

        layout_info = {
            "name": slide_layout.name,
            "placeholders": placeholders_details 
        }
        layout_data["layouts"].append(layout_info)
        print(f"  Analyzed layout: '{slide_layout.name}', Found Placeholders: {list(placeholders_details.keys())}")

    if output_filepath:
        output_dir = os.path.dirname(output_filepath)
        if output_dir:
            os.makedirs(output_dir, exist_ok=True)

        try:
            with open(output_filepath, 'w', encoding='utf-8') as f:
                json.dump(layout_data, f, indent=2)
            print(f"Success! Layouts map saved to: {output_filepath}")
        except Exception as e:
            save_error_message = f"Error saving layouts map to '{output_filepath}': {e}"
            print(save_error_message)
            if "analysis_errors" not in layout_data:
                layout_data["analysis_errors"] = []
            layout_data["analysis_errors"].append(save_error_message)
            
    return layout_data

def main():
    parser = argparse.ArgumentParser(description="Analyze a PowerPoint template and create a layouts.json file.")
    parser.add_argument('-t', '--template', required=True, help="Path to the template PPTX file to analyze.")
    parser.add_argument('-o', '--output', required=True, help="Path for the output layouts.json file.")
    args = parser.parse_args()
    
    analyze_template(args.template, args.output)

if __name__ == '__main__':
    main()
