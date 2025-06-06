from pptx import Presentation

def analyze_presentation(pptx_path):
    prs = Presentation(pptx_path)
    print("\nAnalyzing presentation:", pptx_path)
    
    print("\nAvailable layouts:")
    for layout in prs.slide_layouts:
        print(f"\n{layout.name}:")
        for shape in layout.placeholders:
            print(f"  - {shape.name} (id: {shape.placeholder_format.idx}, type: {shape.placeholder_format.type})")

    print("\nActual slides:")
    for i, slide in enumerate(prs.slides, 1):
        print(f"\nSlide {i}:")
        for shape in slide.placeholders:
            print(f"  - {shape.name} (id: {shape.placeholder_format.idx}, type: {shape.placeholder_format.type})")
            if hasattr(shape, 'text'):
                print(f"    Text: {shape.text[:100] + '...' if len(shape.text) > 100 else shape.text}")

if __name__ == "__main__":
    analyze_presentation("output/my_generated_report.pptx")
