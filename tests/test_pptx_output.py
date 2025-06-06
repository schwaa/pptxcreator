from pptx import Presentation
from pptx.util import Inches, Pt
import json
from collections import defaultdict

class PPTXTester:
    # Standard slide dimensions in EMUs (English Metric Units)
    SLIDE_WIDTH = 9144000  # 10 inches
    SLIDE_HEIGHT = 6858000  # 7.5 inches
    ESTIMATED_LINE_HEIGHT = Pt(20).emu  # ~20pt line height

    def __init__(self, pptx_path, json_path):
        self.pptx_path = pptx_path
        self.json_path = json_path
        self.results = defaultdict(lambda: {"pass": [], "fail": [], "warnings": []})
        
    def load_files(self):
        """Load the PPTX and JSON files."""
        try:
            self.prs = Presentation(self.pptx_path)
            with open(self.json_path, 'r', encoding='utf-8') as f:
                self.data = json.load(f)
            return True
        except Exception as e:
            print(f"Error loading files: {e}")
            return False

    def normalize_text(self, text):
        """Normalize text for comparison."""
        if not text:
            return ""
        # Remove bullets, extra spaces, newlines
        text = text.replace("•", "").replace("\n", " ").replace("  ", " ")
        return text.strip().lower()

    def get_slide_text_by_type(self, slide, placeholder_type):
        """Get text from placeholders of a specific type."""
        text = []
        for shape in slide.placeholders:
            if shape.placeholder_format.type == placeholder_type and hasattr(shape, 'text'):
                text.append(shape.text)
        return '\n'.join(text)

    def get_text_boxes_info(self, slide):
        """Get information about text boxes on the slide."""
        text_boxes = []
        for shape in slide.placeholders:
            if hasattr(shape, 'text'):
                box_info = {
                    'name': shape.name,
                    'type': shape.placeholder_format.type,
                    'text': shape.text,
                    'left': shape.left,
                    'top': shape.top,
                    'width': shape.width,
                    'height': shape.height,
                    'lines': len(shape.text.split('\n'))
                }
                text_boxes.append(box_info)
        return text_boxes

    def check_text_overflow(self, slide, slide_num):
        """Check for text overflow and layout issues."""
        result_key = f"Slide {slide_num}"
        text_boxes = self.get_text_boxes_info(slide)
        
        # Check each text box
        for box in text_boxes:
            # Check for off-slide positioning
            if (box['left'] < 0 or box['top'] < 0 or 
                box['left'] + box['width'] > self.SLIDE_WIDTH or 
                box['top'] + box['height'] > self.SLIDE_HEIGHT):
                self.results[result_key]["fail"].append(
                    f"Shape '{box['name']}' extends beyond slide boundaries")

            # Estimate number of lines that can fit
            max_lines = box['height'] // self.ESTIMATED_LINE_HEIGHT
            actual_lines = box['lines']
            
            if actual_lines > max_lines:
                self.results[result_key]["warnings"].append(
                    f"Possible text overflow in '{box['name']}': "
                    f"Has {actual_lines} lines but space for ~{max_lines}")

        # Check layout usage
        if len(text_boxes) > 1:
            filled_boxes = [box for box in text_boxes if box['text'].strip()]
            if len(filled_boxes) == 1 and len(text_boxes) > 2:
                self.results[result_key]["warnings"].append(
                    "Content concentrated in one box when multiple are available")

    def text_exists_in_slide(self, slide, search_text):
        """Check if text exists anywhere in the slide."""
        if not search_text:
            return True

        # Get all text from placeholders
        slide_text = []
        for shape in slide.placeholders:
            if hasattr(shape, 'text'):
                slide_text.append(shape.text)

        # Normalize and compare
        slide_text = self.normalize_text(' '.join(slide_text))
        search_text = self.normalize_text(search_text)
        
        # Try exact match first
        if search_text in slide_text:
            return True
            
        # Try key phrases
        key_phrases = [p.strip() for p in search_text.split(".") if p.strip()]
        return any(phrase in slide_text for phrase in key_phrases)

    def bullets_exist_in_slide(self, slide, bullets):
        """Check if bullet points exist in slide."""
        if not bullets:
            return True

        # Get text from body placeholders (types 2 and 7)
        slide_text = []
        for shape in slide.placeholders:
            if shape.placeholder_format.type in [2, 7] and hasattr(shape, 'text'):
                slide_text.append(shape.text)
        
        slide_text = self.normalize_text(' '.join(slide_text))
        missing_bullets = []
        
        for bullet in bullets:
            bullet_text = self.normalize_text(bullet)
            if bullet_text not in slide_text:
                missing_bullets.append(bullet)
        
        if missing_bullets:
            self.results[f"Slide {len(self.results)}"]["fail"].append(
                f"Missing bullet points: {missing_bullets}")
            return False
        return True

    def check_slide_content(self, slide_num, slide, content_type, title, **content):
        """Check specific content based on slide type."""
        result_key = f"Slide {slide_num}"
        
        # Check title
        if title:
            if content_type == "presentation_title":
                title_text = self.get_slide_text_by_type(slide, 3)  # CENTER_TITLE
            else:
                title_text = self.get_slide_text_by_type(slide, 1)  # TITLE

            if self.normalize_text(title) in self.normalize_text(title_text):
                self.results[result_key]["pass"].append(f"Title found: '{title}'")
            else:
                self.results[result_key]["fail"].append(f"Title not found: '{title}'")

        # Check for overflow and layout issues
        self.check_text_overflow(slide, slide_num)

        # Content type specific checks
        if content_type == "presentation_title":
            subtitle = content.get("subtitle")
            subtitle_text = self.get_slide_text_by_type(slide, 4)  # SUBTITLE
            if subtitle and self.normalize_text(subtitle) in self.normalize_text(subtitle_text):
                self.results[result_key]["pass"].append("Subtitle found")
            elif subtitle:
                self.results[result_key]["fail"].append(f"Subtitle not found: '{subtitle}'")

        elif content_type == "bullet_points_summary":
            bullets = content.get("bullets", [])
            if self.bullets_exist_in_slide(slide, bullets):
                self.results[result_key]["pass"].append("All bullet points found")

        elif content_type in ["image_and_description", "description_and_image", "product_showcase"]:
            # Check for image
            has_image = False
            for shape in slide.placeholders:
                if shape.placeholder_format.type == 18:  # PICTURE
                    has_image = True
                    break
            
            if has_image:
                self.results[result_key]["pass"].append("Image placeholder found")
            else:
                self.results[result_key]["fail"].append("No image placeholder found")

            # Check text content
            text_content = content.get("text_content") or content.get("caption")
            if text_content and self.text_exists_in_slide(slide, text_content):
                self.results[result_key]["pass"].append("Text content found")
            elif text_content:
                self.results[result_key]["fail"].append("Text content not found")

        elif content_type == "chart_data_slide":
            # Check for chart
            has_chart = False
            for shape in slide.placeholders:
                if shape.placeholder_format.type == 7:  # OBJECT/CHART
                    has_chart = True
                    break
            
            if has_chart:
                self.results[result_key]["pass"].append("Chart placeholder found")
            else:
                self.results[result_key]["fail"].append("No chart placeholder found")

    def check_slides(self):
        """Check each slide against expected content from JSON."""
        if len(self.prs.slides) != len(self.data["sections"]):
            self.results["general"]["fail"].append(
                f"Expected {len(self.data['sections'])} slides, found {len(self.prs.slides)}")

        # Check each slide
        for i, section in enumerate(self.data["sections"]):
            if i >= len(self.prs.slides):
                self.results["general"]["fail"].append(f"Missing slide {i+1}")
                continue
                
            slide = self.prs.slides[i]
            slide_num = i + 1
            content_type = section["content_type"]
            
            print(f"\nChecking slide {slide_num} ({content_type})...")
            
            self.check_slide_content(
                slide_num,
                slide,
                content_type,
                section.get("title"),
                subtitle=section.get("subtitle"),
                bullets=section.get("bullets"),
                text_content=section.get("text_content"),
                caption=section.get("caption")
            )

    def print_summary(self):
        """Print a summary of all test results."""
        print("\n=== PPTX Test Results ===")
        
        total_passes = 0
        total_fails = 0
        total_warnings = 0
        has_issues = False

        # First print general issues
        if self.results["general"]["fail"]:
            has_issues = True
            print("\nGeneral Issues:")
            for msg in self.results["general"]["fail"]:
                print(f"  ✗ {msg}")

        # Print slide-specific results in order
        for i in range(len(self.data["sections"])):
            section = f"Slide {i+1}"
            results = self.results[section]
            if results["pass"] or results["fail"] or results["warnings"]:
                print(f"\n{section}:")
                if results["pass"]:
                    for msg in results["pass"]:
                        print(f"  ✓ {msg}")
                        total_passes += 1
                if results["fail"]:
                    has_issues = True
                    for msg in results["fail"]:
                        print(f"  ✗ {msg}")
                        total_fails += 1
                if results["warnings"]:
                    for msg in results["warnings"]:
                        print(f"  ⚠ {msg}")
                        total_warnings += 1

        # Print final summary
        print(f"\n=== Final Summary ===")
        print(f"Total Passes: {total_passes}")
        print(f"Total Fails: {total_fails}")
        print(f"Total Warnings: {total_warnings}")
        if not has_issues and total_warnings == 0:
            print("\n✓ All checks passed!")
        elif has_issues:
            print("\n✗ Some checks failed. See details above.")
        else:
            print("\n⚠ Tests passed with warnings. See details above.")

def test_pptx_output(pptx_path="output/my_generated_report.pptx", 
                    json_path="data/example_report_data.json"):
    """Run tests on generated PPTX file."""
    print(f"\nTesting PPTX output...")
    print(f"PPTX file: {pptx_path}")
    print(f"JSON file: {json_path}")
    
    tester = PPTXTester(pptx_path, json_path)
    if tester.load_files():
        tester.check_slides()
        tester.print_summary()
    else:
        print("Testing failed: Could not load required files.")

if __name__ == "__main__":
    test_pptx_output()
